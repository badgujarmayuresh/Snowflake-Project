"""
Injects architecture_overview.png as a new slide (slide 2) into project_proposal.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path
import copy

BASE    = Path(__file__).parent.parent
PPT_IN  = BASE / "project_proposal.pptx"
IMG     = BASE / "architecture_overview.png"
PPT_OUT = BASE / "project_proposal.pptx"

DARK_BG  = RGBColor(0x0D, 0x1B, 0x2A)
SNOW_BLUE= RGBColor(0x29, 0xB5, 0xE8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(str(PPT_IN))
BLANK = prs.slide_layouts[6]

# ── Build new slide ────────────────────────────────────────────
new_slide = prs.slides.add_slide(BLANK)

# dark background
def add_rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=WHITE):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

add_rect(new_slide, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(new_slide, 0, 0, 13.33, 0.45, SNOW_BLUE)
add_text(new_slide, "Architecture Overview  —  Single Page Visual", 0, 0.04, 13.33, 0.38,
         size=14, bold=True, color=DARK_BG)

# add the diagram image
new_slide.shapes.add_picture(str(IMG), Inches(0.1), Inches(0.5),
                              Inches(13.13), Inches(6.85))

# ── Move new slide to position 2 (index 1) ───────────────────
xml_slides = prs.slides._sldIdLst
# last added slide is at end — move it to index 1
slide_id_elem = xml_slides[-1]
xml_slides.remove(slide_id_elem)
xml_slides.insert(1, slide_id_elem)

prs.save(str(PPT_OUT))
print(f"Saved updated PPT: {PPT_OUT}")
