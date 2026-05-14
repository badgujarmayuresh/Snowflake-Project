"""
Generates project_proposal.pptx for Smart BI Agent project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from pathlib import Path

OUT = Path(__file__).parent.parent / "project_proposal.pptx"

# ── Brand colours ──────────────────────────────────────────────
SNOW_BLUE    = RGBColor(0x29, 0xB5, 0xE8)   # Snowflake blue
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # Dark navy
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xF2, 0xF4, 0xF8)
ACCENT       = RGBColor(0x00, 0xC4, 0xB4)   # Teal accent
TEXT_DARK    = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bg(slide, color=DARK_BG):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.4, fill=DARK_BG)
    add_rect(slide, 0, 1.35, 13.33, 0.08, fill=SNOW_BLUE)
    add_text(slide, title, 0.4, 0.1, 12, 0.8, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12, 0.5, size=14, color=SNOW_BLUE)

def slide_number(slide, n):
    add_text(slide, str(n), 12.7, 7.1, 0.5, 0.3, size=10, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, DARK_BG)
# accent stripe left
add_rect(sl, 0, 0, 0.12, 7.5, fill=SNOW_BLUE)
# large title
add_text(sl, "Smart BI Agent", 0.4, 1.6, 12, 1.2, size=48, bold=True, color=WHITE)
add_text(sl, "Intelligent Analytics Platform for Retail & Manufacturing",
         0.4, 2.85, 12, 0.7, size=20, color=SNOW_BLUE)
add_rect(sl, 0.4, 3.7, 4.5, 0.06, fill=ACCENT)
add_text(sl, "Prepared by: Mayuresh Prakash Badgujar", 0.4, 3.9, 8, 0.4, size=14, color=LIGHT_GREY)
add_text(sl, "May 2026  |  Version 1.0",               0.4, 4.35, 8, 0.4, size=13, color=RGBColor(0x88,0x99,0xAA))
add_text(sl, "Powered by Snowflake Cortex AI",          0.4, 6.7,  8, 0.4, size=11,
         color=RGBColor(0x55,0x66,0x77), italic=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "Executive Summary", "What are we building and why?")

points = [
    ("The Challenge",
     "Business users cannot self-serve data. Analysts are overwhelmed. Decisions are delayed."),
    ("The Solution",
     "A single AI-powered analytics platform where anyone can ask questions in plain English and get instant answers."),
    ("The Platform",
     "Built entirely on Snowflake — combining a Bronze/Silver/Gold data pipeline with Cortex AI, Cortex Search, and Cortex ML."),
    ("The Outcome",
     "From hours to seconds. From analysts-only to everyone. From reactive to predictive."),
]

for i, (title, body) in enumerate(points):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.7 + row * 2.6
    add_rect(sl, x, y, 6.1, 2.2, fill=WHITE, line=SNOW_BLUE)
    add_rect(sl, x, y, 6.1, 0.45, fill=SNOW_BLUE)
    add_text(sl, title, x+0.15, y+0.05, 5.8, 0.4, size=13, bold=True, color=WHITE)
    add_text(sl, body,  x+0.15, y+0.55, 5.8, 1.5, size=12, color=TEXT_DARK)

slide_number(sl, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "The Problem", "Why does this platform need to exist?")

problems = [
    ("No Self-Service",      "Only analysts can query data.\nBusiness waits days for answers."),
    ("Buried Insights",      "Reports exist but no one can\nfind or search through them."),
    ("No Predictions",       "Teams react to problems.\nNo forward-looking intelligence."),
    ("Fragmented Tools",     "BI + ML + Apps are separate.\nHigh cost, inconsistent data."),
]

for i, (title, body) in enumerate(problems):
    x = 0.35 + i * 3.15
    add_rect(sl, x, 1.7, 2.85, 4.5, fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD))
    add_rect(sl, x, 1.7, 2.85, 0.5, fill=RGBColor(0xE8,0x3A,0x3A))
    add_text(sl, str(i+1), x+0.1, 1.72, 0.4, 0.45, size=16, bold=True, color=WHITE)
    add_text(sl, title, x+0.45, 1.75, 2.3, 0.4, size=13, bold=True, color=WHITE)
    add_text(sl, body,  x+0.15, 2.35, 2.6, 3.5, size=12, color=TEXT_DARK)

slide_number(sl, 3)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — DATA DOMAIN
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "Data Domain", "One company — from factory floor to customer doorstep")

domains = [
    ("E-commerce / Retail",   SNOW_BLUE,
     "Orders  |  Customers  |  Payments\nDelivery  |  Reviews  |  Products"),
    ("B2B / Wholesale",        ACCENT,
     "Suppliers  |  Product Catalogue\nEmployees  |  Wholesale Orders"),
    ("Manufacturing",          RGBColor(0xFF,0x8C,0x00),
     "Machines  |  Production Orders\nDefects  |  Downtime  |  Inventory"),
]

for i, (title, color, body) in enumerate(domains):
    x = 0.5 + i * 4.2
    add_rect(sl, x, 1.7, 3.9, 4.5, fill=WHITE, line=color)
    add_rect(sl, x, 1.7, 3.9, 0.55, fill=color)
    add_text(sl, title, x+0.2, 1.77, 3.6, 0.45, size=14, bold=True, color=WHITE)
    add_text(sl, body,  x+0.2, 2.45, 3.5, 3.5, size=13, color=TEXT_DARK)

add_text(sl, "Combined they represent a company that manufactures and sells products online — end to end.",
         0.5, 6.4, 12.3, 0.5, size=12, italic=True,
         color=RGBColor(0x44,0x55,0x66), align=PP_ALIGN.CENTER)

slide_number(sl, 4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — LAYERED ARCHITECTURE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "Layered Architecture", "Bronze → Silver → Gold → AI → App")

layers = [
    ("SOURCE DATA",  "CSV files: Olist, Northwind, Manufacturing",  RGBColor(0x78,0x78,0x78)),
    ("BRONZE",       "Raw data — exact copy, no transformations",    RGBColor(0xCD,0x7F,0x32)),
    ("SILVER",       "Cleaned, typed, deduplicated, conformed",      RGBColor(0xAA,0xAA,0xAA)),
    ("GOLD",         "Dimensional model — Facts + Dimensions",       RGBColor(0xFF,0xD7,0x00)),
    ("AI LAYER",     "Cortex Analyst + Cortex Search + Cortex ML",  SNOW_BLUE),
    ("STREAMLIT APP","Business user interface — ask anything",       ACCENT),
]

box_w, box_h, gap = 9.5, 0.62, 0.1
start_x, start_y = 1.9, 1.55

for i, (label, desc, color) in enumerate(layers):
    y = start_y + i * (box_h + gap)
    add_rect(sl, start_x, y, box_w, box_h, fill=color)
    lbl_color = TEXT_DARK if color == RGBColor(0xFF,0xD7,0x00) else WHITE
    add_text(sl, label, start_x+0.15, y+0.1, 2.2, box_h-0.15, size=12, bold=True, color=lbl_color)
    add_text(sl, desc,  start_x+2.5,  y+0.1, 6.8, box_h-0.15, size=11, color=lbl_color)
    if i < len(layers)-1:
        add_text(sl, "▼", 6.3, y+box_h, 0.5, gap+0.05, size=9,
                 color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.CENTER)

slide_number(sl, 5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — DIMENSIONAL MODEL
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "Gold Layer — Dimensional Model", "Star Schema design for analytics and AI consumption")

facts = [
    ("FACT_SALES",      "Revenue, qty, freight, discount"),
    ("FACT_ORDERS",     "Order count, payment, delivery days"),
    ("FACT_PRODUCTION", "Planned vs actual qty, cost"),
    ("FACT_DEFECTS",    "Defect qty, defect rate"),
]
dims = [
    ("DIM_DATE",     "Year, Quarter, Month, Week"),
    ("DIM_CUSTOMER", "Customer, city, state, region"),
    ("DIM_PRODUCT",  "Product, category, translation"),
    ("DIM_SELLER",   "Seller details and location"),
    ("DIM_MACHINE",  "Type, plant, status, capacity"),
    ("DIM_LOCATION", "Geographic hierarchy"),
]

# Facts column
add_text(sl, "FACT TABLES", 0.3, 1.55, 4, 0.35, size=12, bold=True, color=SNOW_BLUE)
for i, (name, desc) in enumerate(facts):
    y = 2.0 + i * 1.15
    add_rect(sl, 0.3, y, 4.5, 0.95, fill=RGBColor(0xFF,0xD7,0x00))
    add_text(sl, name, 0.5, y+0.05, 4.1, 0.4, size=12, bold=True, color=TEXT_DARK)
    add_text(sl, desc, 0.5, y+0.45, 4.1, 0.4, size=10, color=TEXT_DARK)

# Dims column
add_text(sl, "DIMENSION TABLES", 5.5, 1.55, 7, 0.35, size=12, bold=True, color=SNOW_BLUE)
for i, (name, desc) in enumerate(dims):
    col = i % 2
    row = i // 2
    x = 5.5 + col * 3.8
    y = 2.0 + row * 1.65
    add_rect(sl, x, y, 3.5, 1.35, fill=SNOW_BLUE)
    add_text(sl, name, x+0.15, y+0.1,  3.2, 0.45, size=11, bold=True, color=WHITE)
    add_text(sl, desc, x+0.15, y+0.6,  3.2, 0.65, size=10, color=LIGHT_GREY)

slide_number(sl, 6)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — AI CAPABILITIES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "AI Capabilities", "Three AI engines working together")

caps = [
    ("Cortex Analyst",  SNOW_BLUE,
     "Natural Language to SQL",
     '"What were top 10 products\nby revenue last quarter?"',
     "Ask questions in plain English.\nGet SQL + results + charts instantly."),
    ("Cortex Search",   ACCENT,
     "Semantic Insight Search",
     '"Find reports about\nsales decline in 2023"',
     "Search across historical reports\nand insights semantically."),
    ("Cortex ML",       RGBColor(0xFF,0x8C,0x00),
     "Forecasting & Anomaly Detection",
     '"Forecast next quarter\ndemand for top products"',
     "Predict future trends.\nDetect anomalies automatically."),
]

for i, (name, color, subtitle, example, desc) in enumerate(caps):
    x = 0.35 + i * 4.3
    add_rect(sl, x, 1.6, 4.0, 5.4, fill=WHITE, line=color)
    add_rect(sl, x, 1.6, 4.0, 0.6, fill=color)
    add_text(sl, name,     x+0.15, 1.65, 3.7, 0.45, size=14, bold=True, color=WHITE)
    add_text(sl, subtitle, x+0.15, 2.35, 3.7, 0.4,  size=11, bold=True, color=color)
    add_rect(sl, x+0.15, 2.85, 3.7, 1.3, fill=LIGHT_GREY)
    add_text(sl, example,  x+0.25, 2.9,  3.5, 1.2,  size=11, italic=True, color=TEXT_DARK)
    add_text(sl, desc,     x+0.15, 4.3,  3.7, 1.5,  size=11, color=TEXT_DARK)

slide_number(sl, 7)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "Technology Stack", "Everything runs natively on Snowflake")

stack = [
    ("Cloud Platform",    "Snowflake",                  "Single platform — storage, compute, AI, app"),
    ("Data Pipeline",     "Bronze / Silver / Gold",      "Structured, layered data architecture"),
    ("Dimensional Model", "Star Schema",                 "Optimised for analytics and AI consumption"),
    ("Natural Language",  "Cortex Analyst",              "Business questions translated to SQL"),
    ("Semantic Search",   "Cortex Search Service",       "Index and retrieve insights semantically"),
    ("Predictive AI",     "Cortex ML Functions",         "Built-in forecasting and anomaly detection"),
    ("Orchestration",     "Snowflake Cortex Agent",      "Intelligent routing across all AI tools"),
    ("User Interface",    "Streamlit in Snowflake",      "Native app — no infrastructure required"),
]

col_headers = ["Layer", "Technology", "Purpose"]
col_x = [0.3, 3.5, 7.0]
col_w = [3.0, 3.3, 6.0]

# header row
for j, (hdr, cx, cw) in enumerate(zip(col_headers, col_x, col_w)):
    add_rect(sl, cx, 1.55, cw-0.1, 0.45, fill=DARK_BG)
    add_text(sl, hdr, cx+0.1, 1.58, cw-0.2, 0.38, size=12, bold=True, color=WHITE)

for i, (layer, tech, purpose) in enumerate(stack):
    y = 2.1 + i * 0.57
    bg = WHITE if i % 2 == 0 else LIGHT_GREY
    for j, (val, cx, cw) in enumerate(zip([layer, tech, purpose], col_x, col_w)):
        add_rect(sl, cx, y, cw-0.1, 0.5, fill=bg, line=RGBColor(0xDD,0xDD,0xDD))
        tcol = SNOW_BLUE if j == 1 else TEXT_DARK
        add_text(sl, val, cx+0.1, y+0.05, cw-0.2, 0.4, size=11,
                 bold=(j==1), color=tcol)

slide_number(sl, 8)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS VALUE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "Business Value", "Measurable impact on speed, access, and intelligence")

metrics = [
    ("Time to Answer",      "Hours / Days",  "Seconds",    SNOW_BLUE),
    ("Self-Service Users",  "Analysts Only", "Everyone",   ACCENT),
    ("Forecast Access",     "Manual/Excel",  "Automated",  RGBColor(0xFF,0x8C,0x00)),
    ("Anomaly Detection",   "Reactive",      "Proactive",  RGBColor(0x8B,0x5C,0xF6)),
    ("Tools Required",      "Multiple",      "One",        RGBColor(0xE8,0x3A,0x3A)),
]

add_text(sl, "BEFORE", 4.5, 1.55, 3.5, 0.4, size=13, bold=True,
         color=RGBColor(0xE8,0x3A,0x3A), align=PP_ALIGN.CENTER)
add_text(sl, "AFTER",  8.2, 1.55, 3.5, 0.4, size=13, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER)

for i, (metric, before, after, color) in enumerate(metrics):
    y = 2.05 + i * 0.95
    add_rect(sl, 0.3, y, 3.8, 0.75, fill=color)
    add_text(sl, metric, 0.45, y+0.12, 3.5, 0.55, size=13, bold=True, color=WHITE)
    add_rect(sl, 4.3, y, 3.6, 0.75, fill=WHITE, line=RGBColor(0xE8,0x3A,0x3A))
    add_text(sl, before, 4.4, y+0.12, 3.4, 0.55, size=13, color=RGBColor(0xE8,0x3A,0x3A),
             align=PP_ALIGN.CENTER)
    add_text(sl, "-->",  8.0, y+0.12, 0.5, 0.55, size=12, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_rect(sl, 8.2, y, 3.6, 0.75, fill=WHITE, line=ACCENT)
    add_text(sl, after,  8.3, y+0.12, 3.4, 0.55, size=13, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)

slide_number(sl, 9)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — DELIVERY PHASES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, LIGHT_GREY)
add_header_bar(sl, "Delivery Phases", "Step-by-step build sequence")

phases = [
    ("1", "Architecture & Design",      "Define layers, data model, tech decisions",         SNOW_BLUE),
    ("2", "Bronze Layer",               "Ingest raw CSV data into Snowflake as-is",           RGBColor(0xCD,0x7F,0x32)),
    ("3", "Silver Layer",               "Clean, standardise and type-cast data",              RGBColor(0xAA,0xAA,0xAA)),
    ("4", "Gold Layer",                 "Build dimensional model — facts + dimensions",       RGBColor(0xFF,0xD7,0x00)),
    ("5", "Cortex Analyst",             "Natural language interface on Gold layer",           ACCENT),
    ("6", "Cortex Search",              "Semantic search across insights",                    RGBColor(0x8B,0x5C,0xF6)),
    ("7", "Cortex ML",                  "Forecasting and anomaly detection",                  RGBColor(0xFF,0x8C,0x00)),
    ("8", "Snowflake Agent",            "Orchestrate all AI tools intelligently",             RGBColor(0xE8,0x3A,0x3A)),
    ("9", "Streamlit App",              "Business user interface — all in one place",         RGBColor(0x06,0xB6,0xD4)),
]

for i, (num, phase, desc, color) in enumerate(phases):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.35
    y = 1.65 + row * 1.85
    add_rect(sl, x, y, 4.1, 1.6, fill=WHITE, line=color)
    add_rect(sl, x, y, 0.55, 1.6, fill=color)
    add_text(sl, num,   x+0.1,  y+0.55, 0.4,  0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    lbl_color = TEXT_DARK if color == RGBColor(0xFF,0xD7,0x00) else color
    add_text(sl, phase, x+0.65, y+0.1,  3.3,  0.5,  size=12, bold=True, color=lbl_color)
    add_text(sl, desc,  x+0.65, y+0.65, 3.3,  0.8,  size=10, color=TEXT_DARK)

slide_number(sl, 10)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl, DARK_BG)
add_rect(sl, 0, 0, 0.12, 7.5, fill=SNOW_BLUE)
add_rect(sl, 0, 3.45, 13.33, 0.08, fill=SNOW_BLUE)

add_text(sl, "Let's Build It.", 0.4, 1.5,  12, 1.2, size=48, bold=True, color=WHITE)
add_text(sl, "Traditional BI expertise meets Snowflake AI.",
         0.4, 2.85, 12, 0.6, size=20, color=SNOW_BLUE)
add_text(sl, "Mayuresh Prakash Badgujar",
         0.4, 4.1,  8,  0.5, size=16, bold=True, color=WHITE)
add_text(sl, "BI & Analytics | 10+ Years Experience | Snowflake Cortex AI",
         0.4, 4.65, 10, 0.5, size=13, color=RGBColor(0x88,0x99,0xAA))
add_text(sl, "Built with Snowflake Cortex AI  |  Streamlit  |  Star Schema Dimensional Modelling",
         0.4, 6.9, 12, 0.4, size=10, color=RGBColor(0x44,0x55,0x66), italic=True)


prs.save(str(OUT))
print(f"Saved: {OUT}")
